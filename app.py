# app.py
# -*- coding: utf-8 -*-
import os
import re
import json
import time
from datetime import datetime
from typing import Dict, List, Tuple, Optional, Any

import requests
from flask import Flask, render_template, request, jsonify, Response
from dotenv import load_dotenv

from docx import Document
from pypdf import PdfReader

from openpyxl import load_workbook
import xlrd

from pptx import Presentation


load_dotenv()

APP_TITLE = "CHUPPY RAG CONVERTER"
HEADER_MODEL_LABEL = "Model : ChatGPT 5.2"

# Chat API (/chat-messages)
API_BASE = (os.getenv("DIFY_API_BASE") or "").strip().rstrip("/")
API_KEY = (os.getenv("DIFY_API_KEY") or "").strip()

# Knowledge API (/datasets ...)
DATASET_API_BASE = (os.getenv("DIFY_DATASET_API_BASE") or API_BASE).strip().rstrip("/")
DATASET_API_KEY = (os.getenv("DIFY_DATASET_API_KEY") or API_KEY).strip()
DATASET_NAME_PREFIX = (os.getenv("DATASET_NAME_PREFIX") or "Chu_").strip()

ALLOWED_EXTS = {
    ".txt", ".md", ".csv", ".json", ".log",
    ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
    ".py", ".js", ".css",
    ".docx", ".pdf",
    ".xlsx", ".xls", ".xlsm",
    ".ppt", ".pptx",
}

MAX_INPUT_CHARS = 180_000
DEFAULT_CHUNK_SEP = "***"
REQ_TIMEOUT_SEC = 300

# Dify knowledge indexing polling
INDEXING_POLL_SEC = float(os.getenv("DIFY_INDEXING_POLL_SEC") or "2.0")
INDEXING_MAX_WAIT_SEC = int(os.getenv("DIFY_INDEXING_MAX_WAIT_SEC") or "900")

# Dify segmentation max_tokens upper bound (per your env)
DIFY_MAX_SEG_TOKENS = int(os.getenv("DIFY_MAX_SEG_TOKENS") or "2000")

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
NOTICE_PATH = os.path.join(BASE_DIR, "notice.txt")


def ensure_notice_file() -> None:
    if os.path.exists(NOTICE_PATH):
        return
    try:
        with open(NOTICE_PATH, "w", encoding="utf-8") as f:
            f.write("")
    except Exception:
        pass


def create_app():
    ensure_notice_file()

    app = Flask(__name__)
    app.config["JSON_AS_ASCII"] = False

    @app.get("/")
    def index():
        return render_template(
            "index.html",
            title=APP_TITLE,
            model_label=HEADER_MODEL_LABEL,
            api_ready=bool(API_BASE and API_KEY),
        )

    @app.get("/auto")
    def auto_page():
        return render_template(
            "auto_rag.html",
            title=APP_TITLE + " (AUTO)",
            model_label=HEADER_MODEL_LABEL,
            api_ready=bool(API_BASE and API_KEY),
            dataset_api_ready=bool(DATASET_API_BASE and DATASET_API_KEY),
            dataset_prefix=DATASET_NAME_PREFIX,
        )

    @app.get("/knowledge")
    def knowledge_page():
        return render_template(
            "knowledge.html",
            title=APP_TITLE + " (KNOWLEDGE)",
            model_label=HEADER_MODEL_LABEL,
            dataset_api_ready=bool(DATASET_API_BASE and DATASET_API_KEY),
            dataset_prefix=DATASET_NAME_PREFIX,
        )

    @app.get("/api/health")
    def api_health():
        return jsonify({
            "ok": True,
            "api_ready": bool(API_BASE and API_KEY),
            "dataset_api_ready": bool(DATASET_API_BASE and DATASET_API_KEY),
            "dataset_prefix": DATASET_NAME_PREFIX,
            "model_label": HEADER_MODEL_LABEL,
        })

    @app.get("/api/notice")
    def api_notice():
        ensure_notice_file()
        try:
            with open(NOTICE_PATH, "r", encoding="utf-8", errors="ignore") as f:
                txt = f.read()
        except Exception:
            return jsonify({"ok": False, "error": "notice.txt の読み取りに失敗しました。"}), 500

        if len(txt) > 50_000:
            txt = txt[:50_000] + "\n...(truncated)\n"

        return jsonify({"ok": True, "text": txt})

    @app.get("/api/datasets")
    def api_datasets():
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            items = dify_list_datasets(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                prefix=DATASET_NAME_PREFIX,
                limit=100,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({"ok": True, "items": items, "prefix": DATASET_NAME_PREFIX})

    @app.get("/api/knowledge/datasets/<dataset_id>/detail")
    def api_knowledge_dataset_detail(dataset_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            ds = dify_get_dataset_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(ds, dict):
            out.update(ds)
        else:
            out["item"] = ds
        return jsonify(out)

    @app.get("/api/knowledge/datasets/<dataset_id>/documents")
    def api_knowledge_documents(dataset_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        keyword = (request.args.get("keyword") or "").strip()

        try:
            items, total = dify_list_documents_all(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                keyword=keyword,
                limit=100,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({"ok": True, "items": items, "total": total, "keyword": keyword})

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>")
    def api_knowledge_document_detail(dataset_id, document_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        metadata = (request.args.get("metadata") or "without").strip() or "without"

        try:
            doc = dify_get_document_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                metadata=metadata,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(doc, dict):
            out.update(doc)
        else:
            out["item"] = doc
        return jsonify(out)

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>/segments")
    def api_knowledge_segments(dataset_id, document_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        keyword = (request.args.get("keyword") or "").strip()
        status = (request.args.get("status") or "").strip()

        try:
            page = int(request.args.get("page") or "1")
            limit = int(request.args.get("limit") or "20")
        except Exception:
            page = 1
            limit = 20

        page = max(1, page)
        limit = max(1, min(100, limit))

        try:
            res = dify_list_segments_page(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                page=page,
                limit=limit,
                keyword=keyword,
                status=status,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({
            "ok": True,
            "items": res.get("items") or [],
            "has_more": bool(res.get("has_more")),
            "total": int(res.get("total") or 0),
            "page": int(res.get("page") or page),
            "limit": int(res.get("limit") or limit),
        })

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>/segments/<segment_id>")
    def api_knowledge_segment_detail(dataset_id, document_id, segment_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            seg = dify_get_segment_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                segment_id=segment_id,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(seg, dict):
            out.update(seg)
        else:
            out["item"] = seg
        return jsonify(out)

    @app.post("/api/scan")
    def api_scan():
        data = request.get_json(force=True) or {}
        in_dir = (data.get("input_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        if not in_dir or not os.path.isdir(in_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400

        files = list_files(in_dir, recursive=recursive)
        return jsonify({"ok": True, "count": len(files), "files": files})

    @app.post("/api/run")
    def api_run():
        # 手動（変換のみ）
        if not API_BASE or not API_KEY:
            return jsonify({
                "ok": False,
                "error": "サーバー側API設定が未完了です。.env に DIFY_API_BASE / DIFY_API_KEY を設定してください。"
            }), 500

        data = request.get_json(force=True) or {}

        input_dir = (data.get("input_dir") or "").strip()
        output_dir = (data.get("output_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        user = (data.get("user") or "rag_converter").strip()
        knowledge_style = (data.get("knowledge_style") or "rag_markdown").strip()
        chunk_sep = (data.get("chunk_sep") or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP

        overwrite = bool(data.get("overwrite", False))

        if not input_dir or not os.path.isdir(input_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400
        if not output_dir:
            return jsonify({"ok": False, "error": "出力フォルダが未指定です。"}), 400

        os.makedirs(output_dir, exist_ok=True)
        files = list_files(input_dir, recursive=recursive)

        def sse():
            yield sse_event("meta", {
                "title": APP_TITLE,
                "mode": "manual",
                "model": HEADER_MODEL_LABEL,
                "total": len(files),
                "overwrite": overwrite,
            })

            ok_count = 0
            ng_count = 0
            skip_count = 0

            for idx, relpath in enumerate(files, start=1):
                abspath = os.path.join(input_dir, relpath)
                yield sse_event("progress", {"index": idx, "total": len(files), "file": relpath})

                try:
                    out_path = make_output_path(output_dir, relpath)

                    if (not overwrite) and os.path.exists(out_path):
                        skip_count += 1
                        yield sse_event("skip_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})
                        continue

                    raw_text, meta = extract_text(abspath, knowledge_style=knowledge_style)

                    if not raw_text.strip():
                        raise RuntimeError("抽出テキストが空でした。")

                    if len(raw_text) > MAX_INPUT_CHARS:
                        raw_text = raw_text[:MAX_INPUT_CHARS] + "\n...(truncated)\n"

                    md_body = convert_via_dify_chat_messages_secure(
                        api_base=API_BASE,
                        api_key=API_KEY,
                        user=user,
                        source_path=relpath,
                        source_meta=meta,
                        text=raw_text,
                        knowledge_style=knowledge_style,
                        chunk_sep=chunk_sep,
                    )

                    md_body = normalize_chunk_sep_lines(md_body, chunk_sep)

                    md_save = attach_source_metadata(md_body, source_relpath=relpath, source_abspath=abspath, source_meta=meta)

                    os.makedirs(os.path.dirname(out_path), exist_ok=True)
                    with open(out_path, "w", encoding="utf-8", newline="\n") as f:
                        f.write(md_save)

                    ok_count += 1
                    yield sse_event("done_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})

                except Exception as e:
                    ng_count += 1
                    yield sse_event("error_one", {"file": relpath, "error": safe_err(str(e))})

            yield sse_event("summary", {
                "ok": ok_count,
                "ng": ng_count,
                "skip": skip_count,
                "total": len(files),
                "overwrite": overwrite,
            })

        return Response(sse(), mimetype="text/event-stream")

    @app.post("/api/auto/run")
    def api_auto_run():
        # 自動（変換→ナレッジ登録）
        if not API_BASE or not API_KEY:
            return jsonify({
                "ok": False,
                "error": "生成AI API設定が未完了です。.env に DIFY_API_BASE / DIFY_API_KEY を設定してください。"
            }), 500

        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({
                "ok": False,
                "error": "ナレッジAPI設定が未完了です。.env に DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY を設定してください。"
            }), 500

        data = request.get_json(force=True) or {}

        input_dir = (data.get("input_dir") or "").strip()
        output_dir = (data.get("output_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        dataset_id = (data.get("dataset_id") or "").strip()

        user = (data.get("user") or "rag_converter").strip()
        knowledge_style = (data.get("knowledge_style") or "rag_markdown").strip()
        chunk_sep = (data.get("chunk_sep") or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP

        overwrite = bool(data.get("overwrite", False))

        if not dataset_id:
            return jsonify({"ok": False, "error": "ナレッジ（dataset_id）が未指定です。"}), 400

        if not input_dir or not os.path.isdir(input_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400
        if not output_dir:
            return jsonify({"ok": False, "error": "出力フォルダが未指定です。"}), 400

        os.makedirs(output_dir, exist_ok=True)
        files = list_files(input_dir, recursive=recursive)

        def sse():
            yield sse_event("meta", {
                "title": APP_TITLE,
                "mode": "auto",
                "model": HEADER_MODEL_LABEL,
                "total": len(files),
                "overwrite": overwrite,
                "dataset_id": dataset_id,
                "chunk_sep": chunk_sep,
            })

            ok_count = 0
            ng_count = 0
            skip_count = 0

            for idx, relpath in enumerate(files, start=1):
                abspath = os.path.join(input_dir, relpath)
                yield sse_event("progress", {"index": idx, "total": len(files), "file": relpath})

                try:
                    out_path = make_output_path(output_dir, relpath)

                    if (not overwrite) and os.path.exists(out_path):
                        skip_count += 1
                        yield sse_event("skip_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})
                        continue

                    raw_text, meta = extract_text(abspath, knowledge_style=knowledge_style)

                    if not raw_text.strip():
                        raise RuntimeError("抽出テキストが空でした。")

                    if len(raw_text) > MAX_INPUT_CHARS:
                        raw_text = raw_text[:MAX_INPUT_CHARS] + "\n...(truncated)\n"

                    md_body = convert_via_dify_chat_messages_secure(
                        api_base=API_BASE,
                        api_key=API_KEY,
                        user=user,
                        source_path=relpath,
                        source_meta=meta,
                        text=raw_text,
                        knowledge_style=knowledge_style,
                        chunk_sep=chunk_sep,
                    )

                    md_body = normalize_chunk_sep_lines(md_body, chunk_sep)

                    md_save = attach_source_metadata(md_body, source_relpath=relpath, source_abspath=abspath, source_meta=meta)

                    os.makedirs(os.path.dirname(out_path), exist_ok=True)
                    with open(out_path, "w", encoding="utf-8", newline="\n") as f:
                        f.write(md_save)

                    yield sse_event("done_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})

                    # --- Dify knowledge register (create-by-text) ---
                    # Difyには front-matter を送らない（チャンク解析/区切りが安定）
                    reg = register_markdown_to_dify(
                        dataset_id=dataset_id,
                        doc_name=os.path.basename(out_path),
                        markdown=md_body,
                        chunk_sep=chunk_sep,
                    )

                    yield sse_event("dataset", {
                        "file": relpath,
                        "doc_id": reg.get("doc_id"),
                        "batch": reg.get("batch"),
                        "chunk_sep": reg.get("chunk_sep"),
                        "chunks": reg.get("chunks"),
                        "chunk_tokens_max": reg.get("chunk_tokens_max"),
                        "dify_max_tokens": reg.get("dify_max_tokens"),
                        "search_method": reg.get("search_method"),
                        "message": "ナレッジ登録 受付",
                    })

                    # poll indexing status (log progress)
                    final = None
                    for prog in iter_indexing_status(
                        dataset_id=dataset_id,
                        batch=reg["batch"],
                        doc_id=reg["doc_id"],
                    ):
                        yield sse_event("dataset_progress", {
                            "file": relpath,
                            "doc_id": reg.get("doc_id"),
                            "batch": reg.get("batch"),
                            "status": prog.get("indexing_status"),
                            "completed_segments": prog.get("completed_segments"),
                            "total_segments": prog.get("total_segments"),
                            "error": prog.get("error"),
                            "terminal": bool(prog.get("terminal")),
                        })
                        if prog.get("terminal"):
                            final = prog
                            break

                    if not final:
                        raise RuntimeError("Dify埋め込みの進捗取得に失敗しました。")

                    if (final.get("indexing_status") or "").lower() != "completed":
                        raise RuntimeError(f"Dify埋め込み失敗: status={final.get('indexing_status')} error={final.get('error')}")

                    if int(final.get("total_segments") or 0) <= 0:
                        raise RuntimeError("Dify側で0セグメントのまま完了しました（separator/max_tokens/text を要確認）。")

                    yield sse_event("dataset_done", {
                        "file": relpath,
                        "doc_id": reg.get("doc_id"),
                        "batch": reg.get("batch"),
                        "status": final.get("indexing_status"),
                        "completed_segments": final.get("completed_segments"),
                        "total_segments": final.get("total_segments"),
                        "message": "ナレッジ登録 完了",
                    })

                    ok_count += 1

                except Exception as e:
                    ng_count += 1
                    yield sse_event("error_one", {"file": relpath, "error": safe_err(str(e))})

            yield sse_event("summary", {
                "ok": ok_count,
                "ng": ng_count,
                "skip": skip_count,
                "total": len(files),
                "overwrite": overwrite,
            })

        return Response(sse(), mimetype="text/event-stream")

    return app


# -----------------------
# File utilities
# -----------------------

def list_files(root_dir: str, recursive: bool = True) -> List[str]:
    results: List[str] = []
    root_dir = os.path.abspath(root_dir)

    if recursive:
        for base, _, files in os.walk(root_dir):
            for name in files:
                ext = os.path.splitext(name)[1].lower()
                if ext in ALLOWED_EXTS:
                    abs_path = os.path.join(base, name)
                    rel = os.path.relpath(abs_path, root_dir)
                    results.append(rel)
    else:
        for name in os.listdir(root_dir):
            abs_path = os.path.join(root_dir, name)
            if os.path.isfile(abs_path):
                ext = os.path.splitext(name)[1].lower()
                if ext in ALLOWED_EXTS:
                    results.append(name)

    results.sort()
    return results


def extract_text(path: str, knowledge_style: str = "rag_markdown") -> Tuple[str, Dict[str, str]]:
    ext = os.path.splitext(path)[1].lower()
    stat = os.stat(path)
    meta = {
        "filename": os.path.basename(path),
        "ext": ext,
        "size_bytes": str(stat.st_size),
        "mtime": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
    }

    if ext in {
        ".txt", ".md", ".csv", ".json", ".log",
        ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
        ".py", ".js", ".css",
    }:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), meta

    if ext == ".docx":
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                parts.append(t)
        return "\n".join(parts), meta

    if ext == ".pdf":
        return extract_pdf_like(path), meta

    if ext in {".xlsx", ".xlsm", ".xls"}:
        if knowledge_style == "rag_natural":
            text = extract_excel_as_markdown_tables(path, ext)
        else:
            text = extract_excel_as_row_records(path, ext)
        return text, meta

    if ext in {".ppt", ".pptx"}:
        return extract_ppt_like(path, ext), meta

    raise RuntimeError(f"未対応の拡張子です: {ext}")


def extract_pdf_like(path: str) -> str:
    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[PAGE {i+1}]\n{txt}")
    return "\n\n".join(parts)


def extract_excel_as_row_records(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_row_records(path)
    return extract_xlsx_like_as_row_records(path)


def extract_xlsx_like_as_row_records(path: str) -> str:
    wb = load_workbook(path, data_only=True, read_only=True)
    out: List[str] = []

    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("[EMPTY]")
            out.append("")
            continue

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_row_records(path: str) -> str:
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("[EMPTY]")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_excel_as_markdown_tables(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_markdown_tables(path)
    return extract_xlsx_like_as_markdown_tables(path)


def extract_xlsx_like_as_markdown_tables(path: str) -> str:
    MAX_ROWS_PER_SHEET = 200
    wb = load_workbook(path, data_only=True, read_only=True)

    out: List[str] = []
    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("(empty)")
            out.append("")
            continue

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= MAX_ROWS_PER_SHEET:
                out.append(f"(… {MAX_ROWS_PER_SHEET}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_markdown_tables(path: str) -> str:
    MAX_ROWS_PER_SHEET = 200
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("(empty)")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= MAX_ROWS_PER_SHEET:
                out.append(f"(… {MAX_ROWS_PER_SHEET}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_ppt_like(path: str, ext: str) -> str:
    try:
        prs = Presentation(path)
    except Exception:
        if ext == ".ppt":
            raise RuntimeError("`.ppt`（旧形式）は python-pptx で直接読めない場合があります。`.pptx` に変換して再実行してください。")
        raise RuntimeError("PowerPointの解析に失敗しました。ファイル破損または形式が想定外です。")

    parts: List[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text.append(t)

        txt = "\n".join(slide_text)
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[SLIDE {i+1}]\n{txt}")

    return "\n\n".join(parts)


def sanitize_header(cell) -> str:
    if cell is None:
        return ""
    s = str(cell).strip()
    s = re.sub(r"\s+", " ", s)
    return s


def normalize_pdf_like_text(s: str) -> str:
    lines = [ln.rstrip() for ln in s.splitlines()]
    out: List[str] = []
    buf = ""

    def flush():
        nonlocal buf
        if buf:
            out.append(buf)
            buf = ""

    for ln in lines:
        t = ln.strip("\u00a0 ").strip()
        if not t:
            flush()
            out.append("")
            continue
        if len(t) == 1:
            buf += t
        else:
            flush()
            out.append(t)
    flush()

    text = "\n".join(out)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def make_output_path(output_dir: str, rel_input_path: str) -> str:
    """
    仕様: yyyymmdd_hhmmss_元ファイル名.md
    + 入力のサブフォルダ構造は維持
    """
    rel_dir = os.path.dirname(rel_input_path)
    base_name = os.path.splitext(os.path.basename(rel_input_path))[0]
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_name = f"{ts}_{base_name}.md"

    # サブフォルダは sanitize しつつ維持
    safe_dir = sanitize_relpath(rel_dir) if rel_dir else ""
    return os.path.join(output_dir, safe_dir, out_name)


def sanitize_relpath(p: str) -> str:
    if not p:
        return ""
    p = p.replace("..", "__")
    p = re.sub(r'[<>:"|?*]', "_", p)
    return p


def sse_event(event: str, data: Dict[str, Any]) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


def safe_err(msg: str) -> str:
    if not msg:
        return "不明なエラー"
    msg = re.sub(r"(app-[A-Za-z0-9_\-]{10,})", "app-***REDACTED***", msg)
    msg = re.sub(r"(Bearer\s+)[A-Za-z0-9_\-\.=]+", r"\1***REDACTED***", msg, flags=re.IGNORECASE)
    msg = re.sub(r"https?://[^\s]+", "[URL_REDACTED]", msg)
    return msg[:700]


# -----------------------
# Chat conversion (/chat-messages)
# -----------------------

def convert_via_dify_chat_messages_secure(
    api_base: str,
    api_key: str,
    user: str,
    source_path: str,
    source_meta: Dict[str, str],
    text: str,
    knowledge_style: str,
    chunk_sep: str,
) -> str:
    url = f"{api_base}/chat-messages"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    instruction = build_rag_instruction(
        source_path=source_path,
        source_meta=source_meta,
        knowledge_style=knowledge_style,
        chunk_sep=chunk_sep,
    )

    query = (
        instruction
        + "\n\n===== SOURCE TEXT BEGIN =====\n"
        + text
        + "\n===== SOURCE TEXT END =====\n"
    )

    payload = {
        "inputs": {},
        "query": query,
        "response_mode": "blocking",
        "user": user,
    }

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=REQ_TIMEOUT_SEC)
    except requests.RequestException:
        raise RuntimeError("API通信に失敗しました（ネットワーク/タイムアウト）。")

    if r.status_code >= 400:
        raise RuntimeError(f"APIエラー（HTTP {r.status_code}）: {safe_err(r.text)}")

    try:
        data = r.json()
    except Exception:
        raise RuntimeError("APIレスポンスの解析に失敗しました。")

    answer = data.get("answer")
    if not answer or not isinstance(answer, str):
        raise RuntimeError("APIレスポンスが想定外です（answerがありません）。")

    return answer.strip() + "\n"


def build_rag_instruction(source_path: str, source_meta: Dict[str, str], knowledge_style: str, chunk_sep: str) -> str:
    meta_lines = "\n".join([f"- {k}: {v}" for k, v in source_meta.items()])
    ext = (source_meta.get("ext") or "").lower()

    first_chunk_rule = f"""
        # 最初のチャンク（必須）
        - 出力の最初のチャンクは必ず「全体構成（目次/分類）」にする。
        - 形式例：
        - 見出し: 「## 全体構成（目次/分類）」
        - 次の1文: 「このチャンクでは文書全体の構成（目次）と分類方針を示す。」
        - 続けて、章立て（大カテゴリ）と、その中で扱う内容の要約を箇条書きで書く。
        - そのチャンクの末尾に必ず「{chunk_sep}」を単独行で置く。
        """

    excel_rules = ""
    if ext in {".xlsx", ".xls", ".xlsm"} and knowledge_style != "rag_natural":
        excel_rules = f"""
        # Excel特別ルール（標準/FAQ用）
        - 入力には [HEADER] と [ROW] が含まれる。
        - 出力は「データ行（[ROW]）1つにつき、必ずチャンク1つ」にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - [ROW]を統合しない。行同士をまとめない。
        """

    if knowledge_style == "rag_natural":
        style_block = f"""
        出力はMarkdownで「RAG向けMarkdown（自然言語）」として整形する。

        # 手順（必須）
        1) まず文書全体の構成を把握し、上位の章立て（大カテゴリ）を作る。
        2) 次に、人間が指示を出すような自然文で各チャンクの目的を宣言してから本文を置く。
        3) チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    elif knowledge_style == "faq":
        style_block = f"""
        出力はMarkdownで、FAQ形式にする。
        - 質問は具体的に、回答は短く「結論→根拠→例」の順にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    else:
        style_block = f"""
        出力はMarkdownで、RAGに最適化したナレッジへ整形する。
        - 文は「主語 + 述語」でできるだけ明確にする。
        - 検索されやすいキーワード（固有名詞/手順名/条件/例外/閾値）を含める。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - 情報を省略しない（重複は統合可）。
        """

    return f"""
        あなたは「社内RAG用ナレッジ整形AI」である。
        入力された文章を、検索精度が最大化するMarkdownへ再構成する。

        # 変換対象ファイル
        - path: {source_path}
        - meta:
        {meta_lines}

        # 絶対ルール
        - 出力は「変換後Markdown本文のみ」とする（前置き/解説/謝罪/注釈は禁止）。
        - 原文が曖昧な場合は「〜である可能性がある」等で補い、捏造しない。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。

        {first_chunk_rule}

        {excel_rules}

        # スタイル
        {style_block}
        """.strip()


# -----------------------
# Markdown metadata + chunk analysis
# -----------------------

def _yaml_quote(v: str) -> str:
    s = "" if v is None else str(v)
    s = s.replace("\\", "\\\\").replace('"', '\\"')
    return f'"{s}"'


def attach_source_metadata(md: str, source_relpath: str, source_abspath: str, source_meta: Dict[str, str]) -> str:
    """Attach source fullpath + meta as YAML front-matter."""
    fm = {
        "source_relpath": source_relpath,
        "source_abspath": os.path.abspath(source_abspath),
        "source_filename": source_meta.get("filename") or os.path.basename(source_abspath),
        "source_ext": source_meta.get("ext") or os.path.splitext(source_abspath)[1].lower(),
        "source_size_bytes": source_meta.get("size_bytes") or "",
        "source_mtime": source_meta.get("mtime") or "",
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }

    lines = ["---"]
    for k, v in fm.items():
        lines.append(f"{k}: {_yaml_quote(v)}")
    lines.append("---")
    lines.append("")

    body = (md or "").lstrip("\ufeff\n\r ")
    return "\n".join(lines) + body


def normalize_chunk_sep_lines(md: str, chunk_sep: str) -> str:
    """Normalize delimiter lines to exactly chunk_sep (trim spaces)."""
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP
    lines = []
    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            lines.append(sep)
        else:
            lines.append(ln.rstrip("\r"))
    out = "\n".join(lines).strip()
    return out + "\n"


def split_chunks(md: str, chunk_sep: str) -> List[str]:
    chunks: List[str] = []
    buf: List[str] = []
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            txt = "\n".join(buf).strip()
            if txt:
                chunks.append(txt)
            buf = []
        else:
            buf.append(ln)

    last = "\n".join(buf).strip()
    if last:
        chunks.append(last)
    return chunks


def estimate_tokens(text: str) -> int:
    """Rough token estimate. Overestimates slightly to reduce accidental re-splitting."""
    if not text:
        return 0
    total = len(text)
    if total <= 0:
        return 0

    ascii_cnt = sum(1 for ch in text if ord(ch) < 128)
    ascii_ratio = ascii_cnt / total

    chars_per_token = 3.0 if ascii_ratio >= 0.60 else 1.6
    est = int(total / chars_per_token) + 1
    return max(1, est)


def analyze_chunks_for_dify(markdown: str, chunk_sep: str) -> Dict[str, Any]:
    chunks = split_chunks(markdown, chunk_sep)
    lens = [estimate_tokens(c) for c in chunks] if chunks else []

    if not lens:
        return {
            "chunks": 0,
            "chunk_tokens_max": 0,
            "dify_max_tokens": min(1000, DIFY_MAX_SEG_TOKENS),
        }

    max_tok = max(lens)

    target = max_tok + 32
    target = max(200, target)
    target = min(DIFY_MAX_SEG_TOKENS, target)

    return {
        "chunks": len(chunks),
        "chunk_tokens_max": max_tok,
        "dify_max_tokens": target,
    }


# -----------------------
# Dify Knowledge API
# -----------------------

def dify_headers(api_key: str) -> Dict[str, str]:
    return {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }


def dify_list_datasets(api_base: str, api_key: str, prefix: str, limit: int = 100) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    page = 1

    while True:
        url = f"{api_base}/datasets?page={page}&limit={limit}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"datasets取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []

        for it in items:
            name = (it.get("name") or "").strip()
            if prefix and not name.startswith(prefix):
                continue
            did = (it.get("id") or "").strip()
            if did and name:
                out.append({"id": did, "name": name})

        has_more = bool(data.get("has_more"))
        if not has_more:
            break

        page += 1
        if page > 200:
            break

    return out


def dify_get_dataset_detail(api_base: str, api_key: str, dataset_id: str) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    # Dify環境/バージョン差で、このエンドポイントが GET を許可しない場合がある（405）。
    # その場合は詳細表示を諦め、最低限の情報だけ返してドキュメント一覧の表示を継続する。
    if r.status_code == 405:
        return {
            "id": dataset_id,
            "name": "",
            "_note": "datasets/{id} が GET 非対応のため、詳細は省略しました。",
        }
    if r.status_code >= 400:
        raise RuntimeError(f"dataset詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_list_documents_all(
    api_base: str,
    api_key: str,
    dataset_id: str,
    keyword: str = "",
    limit: int = 100,
) -> Tuple[List[Dict[str, Any]], int]:
    items_out: List[Dict[str, Any]] = []
    page = 1
    total = 0

    while True:
        qs = f"page={page}&limit={limit}"
        if keyword:
            qs += "&keyword=" + requests.utils.quote(keyword)

        url = f"{api_base}/datasets/{dataset_id}/documents?{qs}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"documents取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []
        total = int(data.get("total") or total or 0)

        for it in items:
            if isinstance(it, dict):
                items_out.append(it)

        if not bool(data.get("has_more")):
            break

        page += 1
        if page > 200:
            break

    return items_out, total


def dify_get_document_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    metadata: str = "without",
) -> Dict[str, Any]:
    meta = metadata.strip() if metadata else "without"
    if meta not in {"all", "only", "without"}:
        meta = "without"

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}?metadata={meta}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"document詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_list_segments_page(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    page: int = 1,
    limit: int = 20,
    keyword: str = "",
    status: str = "",
) -> Dict[str, Any]:
    qs = f"page={page}&limit={limit}"
    if keyword:
        qs += "&keyword=" + requests.utils.quote(keyword)
    if status:
        qs += "&status=" + requests.utils.quote(status)

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments?{qs}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segments取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return {
        "items": data.get("data") or [],
        "has_more": bool(data.get("has_more")),
        "total": int(data.get("total") or 0),
        "page": int(data.get("page") or page),
        "limit": int(data.get("limit") or limit),
    }


def dify_get_segment_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    segment_id: str,
) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments/{segment_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segment詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_create_document_by_text(
    dataset_id: str,
    name: str,
    text: str,
    chunk_sep: str,
    dify_max_tokens: int,
    search_method: str = "hybrid_search",
) -> Tuple[str, str]:
    """Return (doc_id, batch)."""
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/document/create-by-text"

    payload: Dict[str, Any] = {
        "name": name,
        "text": text,
        "indexing_technique": "high_quality",
        "doc_form": "text_model",
        "process_rule": {
            "mode": "custom",
            "rules": {
                "pre_processing_rules": [
                    {"id": "remove_extra_spaces", "enabled": True},
                    {"id": "remove_urls_emails", "enabled": True},
                ],
                "segmentation": {
                    "separator": chunk_sep,
                    "max_tokens": int(dify_max_tokens),
                },
            },
        },
        # ✅ PATCH不要：ここで hybrid_search を指定する
        "retrieval_model": {
            "search_method": search_method,
            "reranking_enable": False,
            "top_k": 5,
            "score_threshold_enabled": False,
        },
    }

    r = requests.post(url, headers=dify_headers(DATASET_API_KEY), json=payload, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"create-by-text 失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    doc = data.get("document") or {}
    doc_id = (doc.get("id") or "").strip()
    batch = (data.get("batch") or "").strip()

    if not doc_id or not batch:
        raise RuntimeError("create-by-text レスポンスが想定外です（document.id / batch がありません）。")

    return doc_id, batch


def dify_get_indexing_status(dataset_id: str, batch: str) -> List[Dict[str, Any]]:
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/documents/{batch}/indexing-status"
    r = requests.get(url, headers={"Authorization": f"Bearer {DATASET_API_KEY}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"indexing-status 取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return data.get("data") or []


def register_markdown_to_dify(dataset_id: str, doc_name: str, markdown: str, chunk_sep: str) -> Dict[str, Any]:
    # ① Markdownのチャンク長を解析 → max_tokens を自動決定（ファイルごと）
    stats = analyze_chunks_for_dify(markdown, chunk_sep)

    # ② separator をアプリ指定値に固定
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    # ③ retrieval_model.search_method = hybrid_search（create-by-textに同梱）
    doc_id, batch = dify_create_document_by_text(
        dataset_id=dataset_id,
        name=doc_name,
        text=markdown,
        chunk_sep=sep,
        dify_max_tokens=int(stats["dify_max_tokens"]),
        search_method="hybrid_search",
    )

    return {
        "doc_id": doc_id,
        "batch": batch,
        "chunk_sep": sep,
        "chunks": stats["chunks"],
        "chunk_tokens_max": stats["chunk_tokens_max"],
        "dify_max_tokens": stats["dify_max_tokens"],
        "search_method": "hybrid_search",
    }


def iter_indexing_status(dataset_id: str, batch: str, doc_id: str):
    start = time.time()
    last_key = None

    while True:
        if time.time() - start > INDEXING_MAX_WAIT_SEC:
            raise RuntimeError("ナレッジ埋め込みがタイムアウトしました。")

        items = dify_get_indexing_status(dataset_id, batch)

        target = None
        for it in items:
            if (it.get("id") or "").strip() == doc_id:
                target = it
                break

        if not target:
            time.sleep(INDEXING_POLL_SEC)
            continue

        st = (target.get("indexing_status") or "").strip()
        completed = int(target.get("completed_segments") or 0)
        total = int(target.get("total_segments") or 0)
        err = target.get("error")

        key = f"{st}:{completed}/{total}:{err}"
        if key != last_key:
            last_key = key
            terminal = st.lower() in {"completed", "error", "failed", "stopped"}
            yield {
                "indexing_status": st,
                "completed_segments": completed,
                "total_segments": total,
                "error": err,
                "terminal": terminal,
            }

            if terminal:
                return

        time.sleep(INDEXING_POLL_SEC)


if __name__ == "__main__":
    app = create_app()
    app.run(host="0.0.0.0", port=5211, debug=False, threaded=True)